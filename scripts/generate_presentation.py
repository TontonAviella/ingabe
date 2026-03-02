#!/usr/bin/env python3
"""Generate Ingabe presentation deck — styled after Microsoft FarmVibes format."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# ── Brand colors ──
DARK_BG = RGBColor(0x1E, 0x1E, 0x2E)
GREEN_ACCENT = RGBColor(0xA3, 0xE6, 0x35)
DARK_GREEN = RGBColor(0x2D, 0x5A, 0x27)
MID_GREEN = RGBColor(0x3A, 0x7D, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
YELLOW_ACCENT = RGBColor(0xE8, 0xD7, 0x4D)
RED_ACCENT = RGBColor(0xC4, 0x28, 0x1B)
TEAL = RGBColor(0x00, 0x96, 0x88)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.font.italic = italic
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_size=14, color=WHITE, bold=False,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(6), bullet=False):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = alignment
    p.space_before = space_before
    if bullet:
        p.level = 0
    return p


def add_green_circle(slide, left, top, size, text, font_size=20):
    """Add a green circle with icon text (like FarmVibes style)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = MID_GREEN
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape


def add_rounded_box(slide, left, top, width, height, text, bg_color=DARK_GREEN,
                    font_size=12, text_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = True
    return shape


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 1: TITLE
    # ════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_BG)

    # Decorative top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN_ACCENT
    bar.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(4), Inches(0.5),
                 "Ingabe by NozaLabs", font_size=14, color=GREEN_ACCENT, bold=True)

    add_text_box(slide, Inches(0.8), Inches(2.2), Inches(11), Inches(2),
                 "Ingabe: AI-Powered GIS\nfor African Agriculture",
                 font_size=44, color=WHITE, bold=True)

    add_text_box(slide, Inches(0.8), Inches(4.8), Inches(8), Inches(1),
                 "Satellite Imagery  |  Soil Analytics  |  Crop Forecasting  |  Climate Monitoring",
                 font_size=18, color=GREEN_ACCENT)

    add_text_box(slide, Inches(0.8), Inches(6.0), Inches(4), Inches(0.5),
                 "gis.nozalabs.rw", font_size=16, color=LIGHT_GRAY)

    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(6), Inches(0.5),
                 "Empowering farmers, agronomists, and organizations across Rwanda",
                 font_size=13, color=LIGHT_GRAY, italic=True)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 2: THE CHALLENGE — Rwanda context
    # ════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1.2),
                 "Rwanda needs to increase production\nand decrease environmental impact",
                 font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # 4 stat boxes
    stats = [
        ("70%", "of Rwanda's workforce\nis in agriculture", GREEN_ACCENT),
        ("0.6 ha", "Average farm size\n— among smallest globally", GREEN_ACCENT),
        ("33%", "of children under 5\nare stunted from malnutrition", YELLOW_ACCENT),
        ("80%", "of farmers lack access\nto modern agri-data tools", YELLOW_ACCENT),
    ]
    for i, (number, desc, accent) in enumerate(stats):
        left = Inches(0.6 + i * 3.2)
        add_text_box(slide, left, Inches(3.0), Inches(2.8), Inches(0.7),
                     number, font_size=42, color=accent, bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(3.9), Inches(2.8), Inches(1.2),
                     desc, font_size=15, color=LIGHT_GRAY,
                     alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.8),
                 "When geospatial data reaches every farmer, we unlock food security\n"
                 "and sustainability across East Africa.",
                 font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 3: DATA-DRIVEN AGRICULTURE — 3 goals
    # ════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                 "Data-driven agriculture", font_size=36, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.6),
                 "Precision agriculture powered by satellite and AI has been shown to:",
                 font_size=20, color=WHITE, bold=True)

    goals = [
        ("Improve yield", "Up to 20% yield increase through\ntimely interventions"),
        ("Reduce cost", "30% reduction in wasted inputs\nwith targeted application"),
        ("Ensure sustainability", "Monitor and reduce GHG emissions\nfrom agricultural practices"),
    ]
    for i, (title, desc) in enumerate(goals):
        top = Inches(2.5 + i * 1.5)
        add_green_circle(slide, Inches(1.2), top, Inches(0.9), title[0], font_size=28)
        add_text_box(slide, Inches(2.5), top + Inches(0.05), Inches(4), Inches(0.5),
                     title, font_size=24, color=WHITE, bold=True)
        add_text_box(slide, Inches(2.5), top + Inches(0.55), Inches(5), Inches(0.6),
                     desc, font_size=14, color=LIGHT_GRAY)

    add_text_box(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.4),
                 "Ingabe brings these capabilities to every farmer in Rwanda — no hardware required.",
                 font_size=13, color=GREEN_ACCENT, italic=True, alignment=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 4: THE CHALLENGE — high cost of manual data
    # ════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(3),
                 "Across Rwanda, the high cost of manual\ndata collection prevents smallholder\n"
                 "farmers from using precision agriculture.",
                 font_size=32, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(1.5),
                 "Ingabe solves this by using free satellite data (Sentinel-2),\n"
                 "open soil datasets (iSDAsoil), and AI — all accessible\n"
                 "from a web browser with no sensors needed.",
                 font_size=20, color=GREEN_ACCENT, alignment=PP_ALIGN.CENTER, bold=True)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 5: INGABE SOLUTION — what it is
    # ════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
                 "Ingabe: Web-Based GIS for Agriculture", font_size=32, color=WHITE, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.6),
                 "A complete geospatial platform accessible from any browser — no installation, no sensors, no cost.",
                 font_size=16, color=LIGHT_GRAY)

    features = [
        ("Satellite Imagery", "Real-time Sentinel-2 indices:\nNDVI, EVI, NDWI, SAVI, NDRE, NDBI\nat 10m resolution, updated every 5 days"),
        ("Soil Analysis", "21 soil properties from iSDAsoil\nat 30m resolution across Africa:\npH, nitrogen, phosphorus, organic carbon..."),
        ("Weather & Climate", "Daily temperature, rainfall, solar radiation\nfrom Copernicus AgERA5 at 11km\n+ real-time Open-Meteo data"),
        ("Yield Forecasting", "DSSAT crop simulation model\nwith Sentinel-2 data assimilation\nfor maize, rice, beans, sorghum, wheat"),
        ("Emissions Monitoring", "EDGAR v8.0 greenhouse gas data:\nCH4, N2O, CO2, NH3\nby agriculture sector"),
        ("Land Cover", "ESRI 10m annual classification:\ncropland, forest, built area, rangeland\nwith automatic zonal statistics"),
    ]
    for i, (title, desc) in enumerate(features):
        col = i % 3
        row = i // 3
        left = Inches(0.5 + col * 4.2)
        top = Inches(2.2 + row * 2.4)
        add_rounded_box(slide, left, top, Inches(3.8), Inches(0.5),
                        title, bg_color=MID_GREEN, font_size=14)
        add_text_box(slide, left + Inches(0.1), top + Inches(0.6), Inches(3.6), Inches(1.5),
                     desc, font_size=12, color=LIGHT_GRAY)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 6: SAGE AI ASSISTANT
    # ════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
                 "Sage: Your AI-Powered GIS Assistant", font_size=32, color=WHITE, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(10), Inches(0.6),
                 "Ask questions in plain language — Sage executes 50+ geospatial tools automatically.",
                 font_size=16, color=GREEN_ACCENT)

    examples = [
        ('"Show me NDVI for Gasabo district"', "Sage calls Sentinel Hub, retrieves\nvegetation indices, displays results on map"),
        ('"What is the soil pH here?"', "Sage queries iSDAsoil at 30m\nresolution, returns full nutrient profile"),
        ('"Create a 5km buffer and\nfind cropland"', "Sage runs buffer geoprocessing,\nadds land cover layer, calculates statistics"),
        ('"Show crop stress alerts"', "Sage analyzes NDVI anomalies,\nidentifies high-risk districts"),
        ('"Forecast maize yield\nfor this field"', "Sage runs DSSAT model with\nsoil + weather + satellite data"),
        ('"Show methane emissions\nfrom agriculture"', "Sage retrieves EDGAR v8.0 data,\ncolors districts by emission levels"),
    ]
    for i, (query, result) in enumerate(examples):
        col = i % 2
        row = i // 2
        left = Inches(0.5 + col * 6.5)
        top = Inches(2.2 + row * 1.6)

        # Query box
        add_rounded_box(slide, left, top, Inches(2.8), Inches(0.9),
                        query, bg_color=RGBColor(0x2A, 0x2A, 0x3E), font_size=11,
                        text_color=GREEN_ACCENT)
        # Arrow
        add_text_box(slide, left + Inches(2.9), top + Inches(0.15), Inches(0.5), Inches(0.5),
                     "->", font_size=20, color=GREEN_ACCENT, bold=True)
        # Result
        add_text_box(slide, left + Inches(3.3), top + Inches(0.05), Inches(3), Inches(0.9),
                     result, font_size=11, color=LIGHT_GRAY)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 7: ARCHITECTURE DIAGRAM
    # ════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                 "Ingabe Platform Architecture", font_size=32, color=WHITE, bold=True)

    # Data Sources column
    add_text_box(slide, Inches(0.3), Inches(1.3), Inches(2.5), Inches(0.4),
                 "DATA SOURCES", font_size=11, color=GREEN_ACCENT, bold=True,
                 alignment=PP_ALIGN.CENTER)
    sources = ["Sentinel-2\n(10m imagery)", "iSDAsoil\n(30m soil)", "Copernicus AgERA5\n(weather)",
               "EDGAR v8.0\n(emissions)", "ESRI LULC\n(land cover)", "NASA POWER\n(climate)"]
    for i, s in enumerate(sources):
        add_rounded_box(slide, Inches(0.2), Inches(1.8 + i * 0.85), Inches(2.5), Inches(0.7),
                        s, bg_color=TEAL, font_size=10)

    # Processing column
    add_text_box(slide, Inches(3.3), Inches(1.3), Inches(3), Inches(0.4),
                 "AI & PROCESSING", font_size=11, color=GREEN_ACCENT, bold=True,
                 alignment=PP_ALIGN.CENTER)
    processing = ["Sentinel Hub API\n(vegetation indices)", "DSSAT Crop Model\n(yield forecasting)",
                  "QGIS Processing\n(35+ algorithms)", "OpenAI LLM\n(Sage AI assistant)",
                  "ML Models\n(anomaly detection)", "PostGIS\n(spatial database)"]
    for i, p in enumerate(processing):
        add_rounded_box(slide, Inches(3.2), Inches(1.8 + i * 0.85), Inches(3), Inches(0.7),
                        p, bg_color=MID_GREEN, font_size=10)

    # Output column
    add_text_box(slide, Inches(6.8), Inches(1.3), Inches(3), Inches(0.4),
                 "USER SERVICES", font_size=11, color=GREEN_ACCENT, bold=True,
                 alignment=PP_ALIGN.CENTER)
    services = ["Interactive Maps\n(MapLibre GL)", "Choropleth Analysis\n(color by value)",
                "Rwanda Dashboard\n(H3 hexagons)", "Crop Health Alerts\n(anomaly detection)",
                "Yield Risk Assessment\n(district-level)", "Share & Embed\n(collaboration)"]
    for i, s in enumerate(services):
        add_rounded_box(slide, Inches(6.7), Inches(1.8 + i * 0.85), Inches(3), Inches(0.7),
                        s, bg_color=DARK_GREEN, font_size=10)

    # Arrows between columns
    for y_off in [Inches(2.1), Inches(3.8), Inches(5.5)]:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2.8), y_off, Inches(0.4), Inches(0.25))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GREEN_ACCENT
        arrow.line.fill.background()

        arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.3), y_off, Inches(0.4), Inches(0.25))
        arrow2.fill.solid()
        arrow2.fill.fore_color.rgb = GREEN_ACCENT
        arrow2.line.fill.background()

    # User box at right
    add_text_box(slide, Inches(10.2), Inches(1.3), Inches(2.8), Inches(0.4),
                 "USERS", font_size=11, color=GREEN_ACCENT, bold=True,
                 alignment=PP_ALIGN.CENTER)
    users = ["Farmers", "Agronomists", "NGOs", "Government\nAgencies", "Researchers", "Extension\nWorkers"]
    for i, u in enumerate(users):
        add_rounded_box(slide, Inches(10.2), Inches(1.8 + i * 0.85), Inches(2.6), Inches(0.7),
                        u, bg_color=RGBColor(0x4A, 0x4A, 0x5A), font_size=10)

    for y_off in [Inches(2.1), Inches(3.8), Inches(5.5)]:
        arrow3 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(9.8), y_off, Inches(0.4), Inches(0.25))
        arrow3.fill.solid()
        arrow3.fill.fore_color.rgb = GREEN_ACCENT
        arrow3.line.fill.background()

    # ════════════════════════════════════════════════════════════════
    # SLIDE 8: SATELLITE IMAGERY — NDVI explained
    # ════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "Satellite-Powered Crop Monitoring", font_size=32, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
                 "Sentinel-2 provides free 10m imagery every 5 days — Ingabe computes 6 vegetation indices automatically.",
                 font_size=15, color=LIGHT_GRAY)

    # NDVI color bar (simulated with rectangles)
    ndvi_ranges = [
        ("< 0.0", "Water", RGBColor(0x1A, 0x5C, 0xB0)),
        ("0.0-0.2", "Bare Soil", RGBColor(0xCC, 0x33, 0x33)),
        ("0.2-0.4", "Sparse", RGBColor(0xE8, 0x8C, 0x30)),
        ("0.4-0.6", "Moderate", RGBColor(0xE8, 0xD7, 0x4D)),
        ("0.6-0.8", "Healthy", RGBColor(0x66, 0xBB, 0x6A)),
        ("> 0.8", "Very Healthy", RGBColor(0x2E, 0x7D, 0x32)),
    ]
    add_text_box(slide, Inches(0.8), Inches(1.8), Inches(3), Inches(0.4),
                 "NDVI — Vegetation Health Scale", font_size=14, color=GREEN_ACCENT, bold=True)
    for i, (rng, label, clr) in enumerate(ndvi_ranges):
        left = Inches(0.8 + i * 2.0)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(2.3), Inches(1.8), Inches(0.6))
        bar.fill.solid()
        bar.fill.fore_color.rgb = clr
        bar.line.fill.background()
        add_text_box(slide, left, Inches(2.35), Inches(1.8), Inches(0.3),
                     rng, font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(2.6), Inches(1.8), Inches(0.3),
                     label, font_size=10, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Other indices table
    indices = [
        ("NDVI", "Overall vegetation health", "0.4 - 0.8"),
        ("EVI", "Dense canopy analysis", "0.3 - 0.7"),
        ("NDWI", "Water / moisture content", "> 0.0"),
        ("SAVI", "Vegetation on bare soil", "0.3 - 0.7"),
        ("NDRE", "Nitrogen / chlorophyll", "> 0.2"),
        ("NDBI", "Built-up area detection", "> 0.0"),
    ]
    add_text_box(slide, Inches(0.8), Inches(3.3), Inches(6), Inches(0.4),
                 "Six Vegetation Indices Computed Automatically", font_size=14, color=GREEN_ACCENT, bold=True)

    # Table header
    header_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.8), Inches(8), Inches(0.45))
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = MID_GREEN
    header_bg.line.fill.background()
    add_text_box(slide, Inches(0.9), Inches(3.82), Inches(1.8), Inches(0.4),
                 "Index", font_size=12, color=WHITE, bold=True)
    add_text_box(slide, Inches(2.8), Inches(3.82), Inches(3.5), Inches(0.4),
                 "What It Measures", font_size=12, color=WHITE, bold=True)
    add_text_box(slide, Inches(6.5), Inches(3.82), Inches(2), Inches(0.4),
                 "Healthy Range", font_size=12, color=WHITE, bold=True)

    for i, (idx, measure, rng) in enumerate(indices):
        top = Inches(4.3 + i * 0.42)
        if i % 2 == 0:
            row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), top - Inches(0.02), Inches(8), Inches(0.42))
            row_bg.fill.solid()
            row_bg.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x3E)
            row_bg.line.fill.background()
        add_text_box(slide, Inches(0.9), top, Inches(1.8), Inches(0.35),
                     idx, font_size=11, color=GREEN_ACCENT, bold=True)
        add_text_box(slide, Inches(2.8), top, Inches(3.5), Inches(0.35),
                     measure, font_size=11, color=LIGHT_GRAY)
        add_text_box(slide, Inches(6.5), top, Inches(2), Inches(0.35),
                     rng, font_size=11, color=WHITE, bold=True)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 9: SOIL ANALYSIS — precision maps
    # ════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
                 "Precision Soil Maps: 21 Properties at 30m", font_size=32, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
                 "iSDAsoil — Africa's highest-resolution soil dataset, powered by machine learning on thousands of samples.",
                 font_size=15, color=LIGHT_GRAY)

    # Three categories
    categories = [
        ("Nutrients", ["Nitrogen (g/kg)", "Phosphorus (ppm)", "Potassium (ppm)",
                       "Calcium", "Magnesium", "Iron", "Sulphur", "Zinc", "Aluminium"]),
        ("Physical Properties", ["Clay Content (%)", "Sand Content (%)", "Silt Content (%)",
                                  "Bulk Density (g/cm3)", "Stone Content (%)", "Bedrock Depth (cm)",
                                  "USDA Texture Class"]),
        ("Chemical Properties", ["pH (0-14)", "Organic Carbon (g/kg)", "Total Carbon (g/kg)",
                                  "Cation Exchange Capacity", "Fertility Classification"]),
    ]
    for col, (cat_name, props) in enumerate(categories):
        left = Inches(0.5 + col * 4.2)
        add_rounded_box(slide, left, Inches(2.0), Inches(3.8), Inches(0.5),
                        cat_name, bg_color=MID_GREEN, font_size=14)
        txBox = slide.shapes.add_textbox(left + Inches(0.1), Inches(2.7), Inches(3.6), Inches(3.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for j, prop in enumerate(props):
            p = tf.add_paragraph() if j > 0 else tf.paragraphs[0]
            p.text = prop
            p.font.size = Pt(12)
            p.font.color.rgb = LIGHT_GRAY
            p.font.name = "Calibri"
            p.space_before = Pt(4)

    # pH scale at bottom
    add_text_box(slide, Inches(0.8), Inches(6.0), Inches(6), Inches(0.4),
                 "Soil pH Interpretation:", font_size=13, color=GREEN_ACCENT, bold=True)
    ph_ranges = [
        ("< 4.5  Extremely Acidic", RGBColor(0xCC, 0x33, 0x33)),
        ("4.5-5.5  Strongly Acidic", RGBColor(0xE8, 0x8C, 0x30)),
        ("5.5-6.5  Good for Most Crops", RGBColor(0x66, 0xBB, 0x6A)),
        ("6.5-7.5  Ideal", RGBColor(0x2E, 0x7D, 0x32)),
        ("7.5+  Alkaline", RGBColor(0xE8, 0xD7, 0x4D)),
    ]
    for i, (label, clr) in enumerate(ph_ranges):
        left = Inches(0.5 + i * 2.5)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(6.4), Inches(2.3), Inches(0.45))
        bar.fill.solid()
        bar.fill.fore_color.rgb = clr
        bar.line.fill.background()
        add_text_box(slide, left, Inches(6.42), Inches(2.3), Inches(0.4),
                     label, font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 10: YIELD FORECASTING — DSSAT pipeline
    # ════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "Crop Yield Forecasting with DSSAT", font_size=32, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
                 "Combining crop simulation, satellite imagery, and soil data for field-level predictions.",
                 font_size=15, color=LIGHT_GRAY)

    # Pipeline steps
    steps = [
        ("1", "Soil Profile", "iSDAsoil 30m\nSaxton & Rawls\npedotransfer"),
        ("2", "Weather Data", "NASA POWER daily\nTemperature, rain,\nsolar radiation"),
        ("3", "Crop Calendar", "Rwanda RAB standards\nSeason A: Sep-Feb\nSeason B: Feb-Jul"),
        ("4", "DSSAT Simulation", "Crop growth model\nBaseline yield\nestimate (t/ha)"),
        ("5", "Satellite Correction", "Sentinel-2 NDVI\nLAI assimilation\nAdjusted forecast"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        left = Inches(0.3 + i * 2.6)
        add_green_circle(slide, left + Inches(0.6), Inches(2.2), Inches(0.7), num, font_size=22)
        add_text_box(slide, left, Inches(3.1), Inches(2.3), Inches(0.4),
                     title, font_size=14, color=GREEN_ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(3.6), Inches(2.3), Inches(1.2),
                     desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Connect arrows
    for i in range(4):
        left = Inches(1.6 + i * 2.6)
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, Inches(2.45), Inches(0.6), Inches(0.2))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GREEN_ACCENT
        arrow.line.fill.background()

    # Supported crops
    add_text_box(slide, Inches(0.8), Inches(5.2), Inches(10), Inches(0.4),
                 "Supported Crops (Rwanda)", font_size=14, color=GREEN_ACCENT, bold=True)

    crops = [("Maize", "Sep-Feb / Feb-Jul"), ("Rice", "Sep-Feb / Feb-Jul"),
             ("Beans", "Sep-Jan / Feb-Jun"), ("Sorghum", "Sep-Feb / Feb-Jul"),
             ("Wheat", "Feb-Jul (marshlands)")]
    for i, (crop, season) in enumerate(crops):
        left = Inches(0.5 + i * 2.5)
        add_rounded_box(slide, left, Inches(5.7), Inches(2.2), Inches(0.4),
                        crop, bg_color=DARK_GREEN, font_size=12)
        add_text_box(slide, left, Inches(6.2), Inches(2.2), Inches(0.4),
                     season, font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 11: RWANDA DASHBOARD
    # ════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "Rwanda Agriculture Dashboard", font_size=32, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
                 "National-level crop monitoring with H3 hexagonal grid analysis across all 30 districts.",
                 font_size=15, color=LIGHT_GRAY)

    # Dashboard components
    components = [
        ("H3 Hexagonal NDVI Map", "Rwanda divided into hexagonal cells,\neach colored by vegetation health.\n"
         "Hover for exact NDVI values.\nRed = stress, Green = healthy."),
        ("District Statistics", "Agricultural parcels count\n"
         "Average NDVI with trend arrows\nYield risk level (Low/Med/High)\nConfidence score percentage"),
        ("NDVI Time Series Chart", "Area chart showing vegetation\nhealth over time with reference lines:\n"
         "0.2 = bare soil threshold (red)\n0.6 = healthy vegetation (green)"),
        ("ML Recommendations", "AI-generated actionable suggestions:\n"
         "'Consider irrigation in high-risk areas'\n'Vegetation declining - investigate drought'\n"
         "Updated from latest satellite data"),
    ]
    for i, (title, desc) in enumerate(components):
        col = i % 2
        row = i // 2
        left = Inches(0.5 + col * 6.3)
        top = Inches(2.0 + row * 2.5)
        add_rounded_box(slide, left, top, Inches(5.8), Inches(0.5),
                        title, bg_color=MID_GREEN, font_size=14)
        add_text_box(slide, left + Inches(0.2), top + Inches(0.6), Inches(5.5), Inches(1.5),
                     desc, font_size=12, color=LIGHT_GRAY)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 12: EMISSIONS & LAND COVER
    # ════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "Environmental Monitoring", font_size=32, color=WHITE, bold=True)

    # Left: Emissions
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5), Inches(0.5),
                 "Greenhouse Gas Emissions (EDGAR v8.0)", font_size=18, color=GREEN_ACCENT, bold=True)

    gases = [
        ("CH4", "Methane", "Livestock, rice paddies, manure"),
        ("N2O", "Nitrous Oxide", "Fertilizer, manure, crop residues"),
        ("CO2", "Carbon Dioxide", "Agricultural soils"),
        ("NH3", "Ammonia", "Fertilizer, manure, burning"),
    ]
    for i, (sym, name, source) in enumerate(gases):
        top = Inches(2.0 + i * 0.9)
        add_rounded_box(slide, Inches(0.8), top, Inches(1), Inches(0.6),
                        sym, bg_color=RGBColor(0xC4, 0x28, 0x1B), font_size=14)
        add_text_box(slide, Inches(2.0), top + Inches(0.02), Inches(2), Inches(0.3),
                     name, font_size=13, color=WHITE, bold=True)
        add_text_box(slide, Inches(2.0), top + Inches(0.32), Inches(3.5), Inches(0.3),
                     source, font_size=11, color=LIGHT_GRAY)

    # Right: Land Cover
    add_text_box(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(0.5),
                 "Land Cover Classification (ESRI 10m)", font_size=18, color=GREEN_ACCENT, bold=True)

    covers = [
        ("Cropland", RGBColor(0xE8, 0xD7, 0x4D), "Active farmland and crop areas"),
        ("Forest", RGBColor(0x39, 0x7D, 0x49), "Tree canopy > 10m height"),
        ("Built Area", RGBColor(0xC4, 0x28, 0x1B), "Urban, roads, infrastructure"),
        ("Rangeland", RGBColor(0xA8, 0xAB, 0x73), "Grasslands and shrublands"),
        ("Water", RGBColor(0x1A, 0x5C, 0xB0), "Rivers, lakes, reservoirs"),
    ]
    for i, (name, clr, desc) in enumerate(covers):
        top = Inches(2.0 + i * 0.9)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), top, Inches(0.6), Inches(0.6))
        bar.fill.solid()
        bar.fill.fore_color.rgb = clr
        bar.line.fill.background()
        add_text_box(slide, Inches(7.8), top + Inches(0.02), Inches(2), Inches(0.3),
                     name, font_size=13, color=WHITE, bold=True)
        add_text_box(slide, Inches(7.8), top + Inches(0.32), Inches(4), Inches(0.3),
                     desc, font_size=11, color=LIGHT_GRAY)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 13: GEOPROCESSING TOOLS
    # ════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "35+ Geoprocessing Tools", font_size=32, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
                 "Powered by QGIS Processing — all accessible through Sage or the visual interface.",
                 font_size=15, color=LIGHT_GRAY)

    tool_categories = [
        ("Vector Geometry", ["Buffer (km)", "Dissolve", "Fix Geometries",
                             "Geometry by Expression", "Reproject Layer", "Field Calculator"]),
        ("Overlay & Joining", ["Clip", "Intersection", "Spatial Join",
                               "Merge Layers", "Join by Location Summary", "Aggregate"]),
        ("Analysis & Grids", ["Zonal Statistics", "Create Grid (hex/rect)",
                              "Statistics by Categories", "Raster Reproject",
                              "DuckDB SQL Queries", "PostGIS Direct Queries"]),
    ]
    for col, (cat, tools) in enumerate(tool_categories):
        left = Inches(0.5 + col * 4.2)
        add_rounded_box(slide, left, Inches(2.0), Inches(3.8), Inches(0.5),
                        cat, bg_color=MID_GREEN, font_size=14)
        txBox = slide.shapes.add_textbox(left + Inches(0.2), Inches(2.7), Inches(3.5), Inches(3.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        for j, tool in enumerate(tools):
            p = tf.add_paragraph() if j > 0 else tf.paragraphs[0]
            p.text = tool
            p.font.size = Pt(13)
            p.font.color.rgb = LIGHT_GRAY
            p.font.name = "Calibri"
            p.space_before = Pt(8)

    add_text_box(slide, Inches(0.5), Inches(6.2), Inches(12), Inches(0.8),
                 'Example: "Create a 5km buffer around Kigali and calculate cropland percentage"\n'
                 "Sage chains: buffer -> land cover overlay -> zonal statistics -> map visualization",
                 font_size=13, color=GREEN_ACCENT, italic=True, alignment=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 14: DATA SOURCES SUMMARY
    # ════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "Open Data Sources — No Cost, No Sensors", font_size=32, color=WHITE, bold=True)

    sources_data = [
        ("Vegetation Indices", "Sentinel-2 via Sentinel Hub", "10-20m", "5 days", "Global"),
        ("Soil Properties", "iSDAsoil (ML predictions)", "30m", "Static (2020)", "Africa"),
        ("Weather / Climate", "Copernicus AgERA5 + Open-Meteo", "~11km", "Daily", "Global"),
        ("Land Cover", "ESRI 10m Annual LULC", "10m", "Annual", "Global"),
        ("Emissions (GHG)", "EDGAR v8.0 (JRC, EU)", "~11km", "Annual", "Global"),
        ("Crop Yield", "DSSAT + Sentinel-2 assimilation", "Per-field", "On demand", "East Africa"),
        ("Admin Boundaries", "Rwanda NISR", "Vector", "Static", "Rwanda"),
    ]

    # Header
    header_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.5))
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = MID_GREEN
    header_bg.line.fill.background()

    cols = [("Data", 0.6, 2.2), ("Source", 2.9, 3.5), ("Resolution", 6.5, 1.5),
            ("Update Freq.", 8.1, 1.5), ("Coverage", 9.7, 1.5)]
    for name, l, w in cols:
        add_text_box(slide, Inches(l), Inches(1.52), Inches(w), Inches(0.4),
                     name, font_size=12, color=WHITE, bold=True)

    for i, row in enumerate(sources_data):
        top = Inches(2.1 + i * 0.65)
        if i % 2 == 0:
            row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), top - Inches(0.05), Inches(12.3), Inches(0.6))
            row_bg.fill.solid()
            row_bg.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x3E)
            row_bg.line.fill.background()
        for j, (_, l, w) in enumerate(cols):
            add_text_box(slide, Inches(l), top, Inches(w), Inches(0.4),
                         row[j], font_size=11,
                         color=GREEN_ACCENT if j == 0 else LIGHT_GRAY,
                         bold=(j == 0))

    add_text_box(slide, Inches(0.5), Inches(6.4), Inches(12), Inches(0.5),
                 "All data sources are free and open — Ingabe eliminates the need for expensive field sensors.",
                 font_size=14, color=GREEN_ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 15: CLOSING
    # ════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.42), prs.slide_width, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN_ACCENT
    bar.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
                 "Ingabe by NozaLabs", font_size=16, color=GREEN_ACCENT, bold=True)

    add_text_box(slide, Inches(2), Inches(2.5), Inches(9), Inches(1.5),
                 "Empowering African Agriculture\nwith AI and Satellite Data",
                 font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2), Inches(4.5), Inches(9), Inches(1),
                 "gis.nozalabs.rw",
                 font_size=24, color=GREEN_ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2), Inches(5.5), Inches(9), Inches(0.5),
                 "No sensors. No installation. Just a browser.",
                 font_size=18, color=LIGHT_GRAY, italic=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.4),
                 "NozaLabs  |  Kigali, Rwanda  |  app.nozalabs.rw",
                 font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # ── Save ──
    out_path = Path(__file__).parent.parent / "docs" / "Ingabe_Presentation.pptx"
    prs.save(str(out_path))
    print(f"Presentation saved to: {out_path}")
    print(f"Slides: {len(prs.slides)}")
    print(f"File size: {out_path.stat().st_size / 1024:.0f} KB")


if __name__ == "__main__":
    build_presentation()
